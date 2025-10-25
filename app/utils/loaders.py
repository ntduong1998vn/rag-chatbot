import mimetypes
import os
from typing import Dict, List


def read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_pdf(path: str) -> str:
    try:
        import pypdf

        reader = pypdf.PdfReader(path)
        return "\n".join([p.extract_text() or "" for p in reader.pages])
    except Exception:
        # fallback unstructured
        from unstructured.partition.pdf import partition_pdf

        elements = partition_pdf(filename=path)
        return "\n".join([e.text for e in elements if hasattr(e, "text") and e.text])


def read_docx(path: str) -> str:
    try:
        import docx

        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        from unstructured.partition.docx import partition_docx

        elements = partition_docx(filename=path)
        return "\n".join([e.text for e in elements if hasattr(e, "text") and e.text])


def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        from unstructured.partition.pptx import partition_pptx

        elements = partition_pptx(filename=path)
        return "\n".join([e.text for e in elements if hasattr(e, "text") and e.text])


def detect_and_extract(path: str) -> str:
    mt, _ = mimetypes.guess_type(path)
    ext = (os.path.splitext(path)[1] or "").lower()
    try:
        if mt == "application/pdf" or ext == ".pdf":
            return read_pdf(path)
        if ext in (".txt", ".md", ".csv", ".json"):
            return read_text_file(path)
        if ext == ".docx":
            return read_docx(path)
        if ext == ".pptx":
            return read_pptx(path)
        # Fallback unstructured (covers many others)
        from unstructured.partition.auto import partition

        elements = partition(filename=path)
        return "\n".join([e.text for e in elements if hasattr(e, "text") and e.text])
    except Exception as e:
        return ""
    except Exception as e:
        return ""
